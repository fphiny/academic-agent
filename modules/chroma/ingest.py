from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from typing import Any, BinaryIO, Dict, Optional
from xml.etree import ElementTree as ET

import fitz  # PyMuPDF
import pymupdf4llm
from docx import Document
from pptx import Presentation

from .chunking import get_chunker
from .store import get_store

try:
    from modules.agent_speed.tools.external_extract_main_content.service import HTMLPreprocessor
except Exception:
    HTMLPreprocessor = None  # type: ignore


def _get_file_obj(file_storage: Any) -> BinaryIO:
    """
    FastAPI UploadFile -> .file
    Werkzeug FileStorage -> .stream
    둘 다 지원
    """
    if hasattr(file_storage, "file") and file_storage.file is not None:
        return file_storage.file
    if hasattr(file_storage, "stream") and file_storage.stream is not None:
        return file_storage.stream
    raise ValueError("지원하지 않는 업로드 파일 객체입니다.")


def _seek_start(file_storage: Any) -> BinaryIO:
    f = _get_file_obj(file_storage)
    f.seek(0)
    return f


def _read_all(file_storage: Any) -> bytes:
    """
    UploadFile / FileStorage 둘 다 안전하게 전체 bytes 읽기
    """
    f = _seek_start(file_storage)
    raw = f.read()
    if isinstance(raw, str):
        raw = raw.encode("utf-8")
    return raw or b""


def decode_uploaded_text(raw: bytes) -> str:
    encodings = [
        "utf-8",
        "utf-8-sig",
        "cp949",
        "euc-kr",
    ]

    for enc in encodings:
        try:
            return raw.decode(enc).strip()
        except UnicodeDecodeError:
            continue

    return raw.decode("utf-8", errors="replace").strip()


def normalize_text(text: str) -> str:
    text = (text or "").replace("\x00", " ")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text_from_plain(file_storage) -> str:
    raw = _read_all(file_storage)
    if not raw:
        return ""
    return normalize_text(decode_uploaded_text(raw))


def extract_text_from_pdf(file_storage) -> str:
    raw = _read_all(file_storage)
    if not raw:
        return ""

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name

    try:
        try:
            md_text = pymupdf4llm.to_markdown(tmp_path)
            md_text = normalize_text(md_text)
            if md_text:
                return md_text
        except Exception:
            pass

        doc = fitz.open(stream=raw, filetype="pdf")
        pages = []
        for page_idx, page in enumerate(doc, start=1):
            txt = page.get_text("text") or ""
            txt = normalize_text(txt)
            if txt:
                pages.append(f"[Page {page_idx}]\n{txt}")
        return "\n\n".join(pages).strip()
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def extract_text_from_pptx(file_storage) -> str:
    file_obj = _seek_start(file_storage)
    prs = Presentation(file_obj)

    slides_text = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = normalize_text(shape.text)
                if txt:
                    parts.append(txt)

            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_cells = []
                    for cell in row.cells:
                        cell_text = normalize_text(cell.text)
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        parts.append(" | ".join(row_cells))

        if parts:
            slides_text.append(f"[Slide {slide_idx}]\n" + "\n".join(parts))

    return "\n\n".join(slides_text).strip()


def extract_text_from_docx(file_storage) -> str:
    file_obj = _seek_start(file_storage)
    doc = Document(file_obj)

    parts = []

    for para in doc.paragraphs:
        txt = normalize_text(para.text)
        if txt:
            parts.append(txt)

    for table_idx, table in enumerate(doc.tables, start=1):
        table_rows = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = normalize_text(cell.text)
                if cell_text:
                    row_cells.append(cell_text)
            if row_cells:
                table_rows.append(" | ".join(row_cells))
        if table_rows:
            parts.append(f"[Table {table_idx}]")
            parts.extend(table_rows)

    return "\n\n".join(parts).strip()


def extract_text_from_hwpx(file_storage) -> str:
    raw = _read_all(file_storage)
    if not raw:
        return ""

    texts = []

    with tempfile.NamedTemporaryFile(suffix=".hwpx", delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name

    try:
        with zipfile.ZipFile(tmp_path, "r") as zf:
            xml_names = [
                name for name in zf.namelist()
                if name.startswith("Contents/section") and name.endswith(".xml")
            ]

            for name in sorted(xml_names):
                xml_bytes = zf.read(name)
                root = ET.fromstring(xml_bytes)

                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())

        return normalize_text("\n".join(texts))
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def extract_text_from_hwp(file_storage) -> str:
    if shutil.which("hwp5txt") is None:
        raise ValueError(
            "HWP 지원을 위해 hwp5txt 명령이 필요합니다. pyhwp 설치 후 CLI 사용 가능 여부를 확인하세요."
        )

    raw = _read_all(file_storage)
    if not raw:
        return ""

    with tempfile.NamedTemporaryFile(suffix=".hwp", delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name

    try:
        proc = subprocess.run(
            ["hwp5txt", tmp_path],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            check=False,
        )

        if proc.returncode != 0:
            stderr = (proc.stderr or "").strip()
            raise ValueError(f"HWP 텍스트 추출 실패: {stderr or 'unknown error'}")

        return normalize_text(proc.stdout or "")
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def extract_text_by_extension(file_storage, ext: str) -> str:
    if ext in {".txt", ".md", ".csv", ".json"}:
        return extract_text_from_plain(file_storage)

    if ext == ".pdf":
        return extract_text_from_pdf(file_storage)

    if ext == ".pptx":
        return extract_text_from_pptx(file_storage)

    if ext == ".docx":
        return extract_text_from_docx(file_storage)

    if ext == ".hwp":
        return extract_text_from_hwp(file_storage)

    if ext == ".hwpx":
        return extract_text_from_hwpx(file_storage)

    raise ValueError("현재는 txt, md, csv, json, pdf, pptx, docx, hwp, hwpx 파일만 지원합니다.")


def sanitize_metadata_for_vectorstore(data: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Chroma metadata 제약 때문에 list/dict/None 등을 안전한 scalar/string 로 바꾼다.
    """
    cleaned: Dict[str, Any] = {}

    for key, value in (data or {}).items():
        if value is None:
            continue

        if isinstance(value, bool):
            cleaned[key] = value
            continue

        if isinstance(value, (int, float)):
            cleaned[key] = value
            continue

        if isinstance(value, str):
            value = value.strip()
            if value:
                cleaned[key] = value
            continue

        if isinstance(value, list):
            items = [str(v).strip() for v in value if v is not None and str(v).strip()]
            if items:
                cleaned[key] = ", ".join(items)
            continue

        if isinstance(value, dict):
            if value:
                cleaned[key] = json.dumps(value, ensure_ascii=False)
            continue

        text = str(value).strip()
        if text:
            cleaned[key] = text

    return cleaned


def _wrap_html_fragment(raw_html: str) -> str:
    """
    selector 로 잘린 fragment 만 들어와도 preprocessor 가 안정적으로 돌도록 감싼다.
    """
    html = (raw_html or "").strip()
    if not html:
        return ""

    lowered = html.lower()
    if "<html" in lowered or "<body" in lowered:
        return html

    return (
        "<!DOCTYPE html>\n"
        "<html lang=\"ko\">\n"
        "<head><meta charset=\"utf-8\"></head>\n"
        f"<body>\n{html}\n</body>\n"
        "</html>"
    )


def _get_html_preprocessor() -> HTMLPreprocessor:
    if HTMLPreprocessor is None:
        raise ImportError(
            "modules.agent_speed.tools.external_extract_main_content.service.HTMLPreprocessor "
            "를 import 할 수 없습니다."
        )
    return HTMLPreprocessor()


def ingest_text(
    text: str,
    collection_name: str,
    doc_id: Optional[str] = None,
    metadata: Optional[Dict] = None,
):
    chunker = get_chunker()
    store = get_store()

    clean_text = normalize_text(text or "")
    if not clean_text:
        raise ValueError("텍스트가 비어 있습니다.")

    clean_metadata = sanitize_metadata_for_vectorstore(metadata)

    chunks = chunker.chunk_document(
        text=clean_text,
        doc_id=doc_id,
        metadata=clean_metadata,
    )

    ids, docs, metas = chunker.chunks_to_lists(chunks)

    if not ids:
        raise ValueError("청크가 생성되지 않았습니다. 문서 길이나 chunk 설정을 확인하세요.")

    store.upsert_documents(
        collection_name=collection_name,
        ids=ids,
        documents=docs,
        metadatas=metas,
    )

    return {
        "doc_id": doc_id,
        "chunks": len(ids),
    }


def ingest_html_semantic(
    raw_html: str | bytes,
    url: str,
    collection_name: str,
    doc_id: Optional[str] = None,
    title: str = "",
    metadata: Optional[Dict[str, Any]] = None,
):
    """
    웹 선택영역/본문 HTML 전용.
    external_extract_main_content.HTMLPreprocessor 기반 semantic chunking 후 저장한다.
    """
    store = get_store()
    preprocessor = _get_html_preprocessor()

    if isinstance(raw_html, bytes):
        html_input = raw_html
    else:
        html_input = _wrap_html_fragment(raw_html)

    if not html_input:
        raise ValueError("HTML이 비어 있습니다.")

    base_doc_id = (doc_id or "web_document").strip()
    base_url = (url or "").strip()

    base_metadata = sanitize_metadata_for_vectorstore({
        **(metadata or {}),
        "doc_id": base_doc_id,
        "url": base_url,
        "title": title,
        "chunking_strategy": "external_extract_main_content",
    })

    result = preprocessor.process(
        raw_html=html_input,
        url=base_url,
        title=title or "",
    )

    if not result.ok:
        raise ValueError(result.error or "HTML semantic preprocessing failed")

    if not result.chunks:
        raise ValueError("semantic chunk result is empty")

    ids = []
    docs = []
    metas = []
    chunk_items = []

    resolved_title = (result.title or title or "").strip()

    for idx, chunk in enumerate(result.chunks):
        chunk_doc_id = f"{base_doc_id}__{idx:04d}"
        chunk_text = normalize_text(chunk.text or "")
        if not chunk_text:
            continue

        block_type = getattr(chunk.block_type, "value", str(chunk.block_type))

        chunk_meta = sanitize_metadata_for_vectorstore({
            **base_metadata,
            "doc_id": base_doc_id,
            "chunk_doc_id": chunk_doc_id,
            "semantic_chunk_id": chunk.chunk_id,
            "url": (chunk.url or base_url).strip(),
            "title": (chunk.title or resolved_title).strip(),
            "heading": (chunk.heading or "").strip(),
            "block_type": block_type,
            "char_count": int(getattr(chunk, "char_count", len(chunk_text))),
            "source": base_metadata.get("source", "web_scrape_semantic"),
        })

        ids.append(chunk_doc_id)
        docs.append(chunk_text)
        metas.append(chunk_meta)

        chunk_items.append({
            "chunk_doc_id": chunk_doc_id,
            "semantic_chunk_id": chunk.chunk_id,
            "heading": chunk.heading,
            "block_type": block_type,
            "char_count": int(getattr(chunk, "char_count", len(chunk_text))),
        })

    if not ids:
        raise ValueError("semantic chunk 결과가 모두 비어 있습니다.")

    store.upsert_documents(
        collection_name=collection_name,
        ids=ids,
        documents=docs,
        metadatas=metas,
    )

    return {
        "doc_id": base_doc_id,
        "title": resolved_title,
        "chunks": len(ids),
        "total_chars": int(getattr(result, "total_chars", sum(len(x) for x in docs))),
        "items": chunk_items,
        "meta": getattr(result, "meta", {}) or {},
    }


def ingest_file(
    file_storage,
    collection_name: str,
    doc_id: str | None = None,
    metadata: dict | None = None,
):
    filename = getattr(file_storage, "filename", "") or ""
    ext = os.path.splitext(filename)[1].lower()

    if not ext:
        raise ValueError("파일 확장자를 확인할 수 없습니다.")

    text = extract_text_by_extension(file_storage, ext)

    if not text:
        raise ValueError("파일에서 텍스트를 추출하지 못했습니다.")

    if not doc_id:
        doc_id = os.path.splitext(filename)[0] or "uploaded_file"

    merged_metadata = {
        **(metadata or {}),
        "source_file": filename,
        "source_ext": ext,
    }

    return ingest_text(
        text=text,
        collection_name=collection_name,
        doc_id=doc_id,
        metadata=merged_metadata,
    )
